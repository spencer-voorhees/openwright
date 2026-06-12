"""Image-per-slide PPTX exporter — pixel-perfect, not editable.

For each <section> in the deck-stage, screenshots the rendered slide
at the authored 1920x1080 size and embeds it as a full-bleed image in
a pptx. Output is visually identical to the live HTML but the slides
are images, not editable shapes — the trade is pixel fidelity vs
editability. Use this when the deck has visual elements the
DOM-walking exporter can't capture (gradients, complex layouts, custom
fonts the recipient doesn't have).

Usage:
    python3 export_pptx_image.py --url <preview-url> --out <pptx>
"""
from __future__ import annotations
import argparse
import io
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

PX_TO_EMU = 9525   # 1 CSS px @ 96dpi


def main() -> int:
    ap = argparse.ArgumentParser(prog="export_pptx_image",
        description="Image-per-slide HTML → PPTX exporter (pixel-perfect)")
    ap.add_argument("--url", required=True)
    ap.add_argument("--out", required=True)
    ap.add_argument("--width", type=int, default=1920)
    ap.add_argument("--height", type=int, default=1080)
    args = ap.parse_args()

    out_path = Path(args.out).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    from playwright.sync_api import sync_playwright

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(
            viewport={"width": args.width, "height": args.height},
            device_scale_factor=2,    # render at 2x for crispness
        )
        page.goto(args.url, wait_until="networkidle", timeout=30000)
        page.wait_for_function("customElements.get('deck-stage') !== undefined", timeout=15000)
        page.wait_for_timeout(500)

        # Discover how many sections we have. Also nuke all animations
        # + transitions globally so the screenshot captures the final
        # state, not the from-keyframe (opacity:0 + translateY) state
        # the design systems use for their "rise" entry motion.
        n_slides = page.evaluate("""
          () => {
            const stage = document.querySelector('deck-stage');
            if (stage) stage.setAttribute('noscale', '');
            // Inject a style that disables animations + transitions
            // everywhere. The design systems use `animation: ...rise
            // both` which leaves content in the from-state (faded)
            // until the keyframe runs. We want the to-state for
            // capture.
            const s = document.createElement('style');
            s.textContent = `
              *, *::before, *::after {
                animation-duration: 0s !important;
                animation-delay: 0s !important;
                animation-fill-mode: forwards !important;
                animation-iteration-count: 1 !important;
                animation-play-state: running !important;
                transition: none !important;
              }
            `;
            document.head.appendChild(s);
            return document.querySelectorAll('deck-stage > section').length;
          }
        """)
        if not n_slides:
            print("no slides found", file=sys.stderr)
            browser.close()
            return 1

        # Screenshot one slide at a time. Hide all other sections so
        # only the target is visible at (0,0). This sidesteps
        # deck-shell's absolute stacking — whichever slide is
        # display:block renders at the top of the viewport.
        with tempfile.TemporaryDirectory() as td:
            scratch = Path(td)
            shots = []
            for i in range(n_slides):
                page.evaluate(f"""
                  () => {{
                    const sections = Array.from(document.querySelectorAll('deck-stage > section'));
                    sections.forEach((s, idx) => {{
                      if (idx === {i}) {{
                        // Preserve the section's intrinsic display
                        // (flex/grid/block per sample.html). Override
                        // visibility + position + dims only — don't
                        // touch display, or the section's child layout
                        // (padding, gap, alignment) gets nuked.
                        const cur = getComputedStyle(s).display;
                        s.style.cssText = `
                          visibility: visible !important;
                          opacity: 1 !important;
                          position: absolute !important;
                          left: 0 !important;
                          top: 0 !important;
                          width: {args.width}px !important;
                          height: {args.height}px !important;
                          z-index: 9999 !important;
                          ${{cur && cur !== 'none' ? 'display: ' + cur + ' !important;' : 'display: flex !important;'}}
                        `;
                      }} else {{
                        s.style.cssText = 'display: none !important;';
                      }}
                    }});
                    window.scrollTo(0, 0);
                    void document.body.offsetHeight;
                  }}
                """)
                # Let the browser settle after the visibility change.
                page.wait_for_timeout(120)
                png_path = scratch / f"slide-{i + 1}.png"
                page.screenshot(
                    path=str(png_path),
                    clip={"x": 0, "y": 0, "width": args.width, "height": args.height},
                )
                shots.append(png_path)

            browser.close()

            # Assemble pptx: 1 slide per PNG, full-bleed.
            prs = Presentation()
            prs.slide_width  = Emu(args.width  * PX_TO_EMU)
            prs.slide_height = Emu(args.height * PX_TO_EMU)
            blank = prs.slide_layouts[6]
            for png in shots:
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(
                    str(png),
                    0, 0,
                    width=prs.slide_width, height=prs.slide_height,
                )
            prs.save(str(out_path))

    print(f"wrote {out_path}  ({n_slides} slide(s))")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
