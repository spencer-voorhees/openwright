"""DOM-walking HTML → PPTX exporter.

Opens a deck.html in headless chromium at the deck's authored size
(default 1920×1080), inspects each <section> child of <deck-stage>,
captures every visible leaf element's bounds + computed styles +
content, and emits a python-pptx slide where each leaf becomes a real
editable shape (text box, picture, or rect).

Usage:
    python3 pptx.py --url http://host/preview/.../deck.html \\
                     --out  /path/to/deck.pptx

The script reads a JSON primitive manifest from JS executed inside the
page, so the page's own CSS + layout determine where everything sits.
Output pptx is editable — text, images, and shapes are real pptx
primitives, not rasterized. Inline SVGs (icon sets like lucide) are
the one exception: they rasterize to small PNG pictures via canvas,
which keeps them movable/resizable in PowerPoint.
"""
from __future__ import annotations
import argparse
import base64
import io
import json
import re
import sys
import tempfile
import urllib.parse
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# 1 px (CSS px at 96 dpi) = 9525 EMU
PX_TO_EMU = 9525

# JS that runs inside the page. Returns a list of slide primitive sets:
#   [{ width, height, rect, primitives: [ {kind, ...}, ... ] }, ...]
# The IIFE is async: inline SVGs rasterize through an <img> + canvas
# round-trip, which is promise-based. page.evaluate awaits it.
DOM_WALK_JS = r"""
(async () => {
  // Set noscale so the DOM is at authored 1920x1080 coordinates.
  const stage = document.querySelector('deck-stage');
  if (stage) stage.setAttribute('noscale', '');

  // Kill the design systems' entry animations. They animate from
  // opacity:0 — and toggling each section's display restarts the
  // animation from frame 0, so when we re-measure children for
  // slides 1+ they're still opacity:0 and isVisible() rejects them.
  // Match the image exporter: zero animation duration globally.
  const killAnim = document.createElement('style');
  killAnim.textContent = `
    *, *::before, *::after {
      animation-duration: 0s !important;
      animation-delay: 0s !important;
      animation-fill-mode: forwards !important;
      animation-iteration-count: 1 !important;
      animation-play-state: running !important;
      transition: none !important;
    }
  `;
  document.head.appendChild(killAnim);

  const stageRect = stage ? stage.getBoundingClientRect() : { left: 0, top: 0, width: 1920, height: 1080 };
  const sections = Array.from(document.querySelectorAll('deck-stage > section'));

  // The deck-shell hides non-current slides via display: none /
  // visibility: hidden so only one renders at a time. For export we
  // walk one section at a time: show JUST that section (keep its
  // original display: the sample.html relies on display: flex on
  // sections for the layout), measure children, repeat.
  //
  // Crucially: do NOT override `display` — the sections need their
  // original flex layout to position children correctly. Only
  // override visibility/position/dimensions.
  const showOnly = (idx) => {
    sections.forEach((s, i) => {
      if (i === idx) {
        // Preserve original computed display by storing + restoring.
        const cur = getComputedStyle(s).display;
        s.style.cssText = `
          visibility: visible !important;
          opacity: 1 !important;
          position: absolute !important;
          left: 0 !important;
          top: 0 !important;
          width: ${stageRect.width || 1920}px !important;
          height: ${stageRect.height || 1080}px !important;
          z-index: 9999 !important;
          ${cur && cur !== 'none' ? `display: ${cur} !important;` : 'display: flex !important;'}
        `;
      } else {
        s.style.cssText = 'display: none !important;';
      }
    });
    // Force layout pass.
    void document.body.offsetHeight;
  };

  const parseColor = (s) => {
    if (!s || s === 'transparent') return null;
    const m = s.match(/^rgba?\(([^)]+)\)$/i);
    if (!m) return null;
    const parts = m[1].split(',').map((x) => parseFloat(x.trim()));
    const [r, g, b, a] = [parts[0], parts[1], parts[2], parts[3] ?? 1];
    if (a === 0) return null;
    return { r: Math.round(r), g: Math.round(g), b: Math.round(b), alpha: a };
  };

  const toHex = (rgb) =>
    ((rgb.r << 16) | (rgb.g << 8) | rgb.b).toString(16).padStart(6, '0').toUpperCase();

  // Walk up the DOM finding the nearest ancestor with an opaque (alpha=1)
  // background. Used to blend translucent rgba() fills against what the
  // user actually sees, since pptx rects don't support transparency on
  // the basic shape fill.
  const effectiveBg = (el) => {
    let cur = el.parentElement;
    while (cur) {
      const c = parseColor(getComputedStyle(cur).backgroundColor);
      if (c && c.alpha >= 0.999) return c;
      cur = cur.parentElement;
    }
    return { r: 255, g: 255, b: 255, alpha: 1 };   // body fallback
  };

  // Composite a translucent rgba over an opaque base. Returns the
  // visible color the viewer perceives, as a hex string.
  const blendOverBase = (top, base) => {
    if (top.alpha >= 0.999) return toHex(top);
    const a = top.alpha;
    return toHex({
      r: Math.round(top.r * a + base.r * (1 - a)),
      g: Math.round(top.g * a + base.g * (1 - a)),
      b: Math.round(top.b * a + base.b * (1 - a)),
    });
  };

  const isVisible = (el, st) => {
    if (st.display === 'none' || st.visibility === 'hidden' || parseFloat(st.opacity) === 0) return false;
    const r = el.getBoundingClientRect();
    return r.width > 0.5 && r.height > 0.5;
  };

  const hasOnlyTextOrInline = (el) => {
    // True if the element's children are all text nodes or true-inline
    // elements (b, em, span) — i.e., this element owns the text it
    // contains and we should emit one text shape for it. inline-flex /
    // inline-block / inline-grid children do NOT qualify: they paint
    // their own boxes (rating dots, pills) and need the container path
    // or their backgrounds vanish behind the text capture.
    for (const c of el.childNodes) {
      if (c.nodeType === Node.TEXT_NODE) continue;
      if (c.nodeType === Node.ELEMENT_NODE) {
        // An inline <svg> (icon) is NOT text content — handle the
        // element as a container so the svg gets rasterized.
        if (c.tagName && c.tagName.toLowerCase() === 'svg') return false;
        const cs = window.getComputedStyle(c);
        const d = cs.display;
        if (d === 'inline') continue;
        // display:contents is a transparent wrapper — its children
        // hoist into this element's layout, so look through it: a
        // contents row full of block cells must NOT read as inline.
        if (d === 'contents') {
          if (!hasOnlyTextOrInline(c)) return false;
          continue;
        }
        return false;
      }
    }
    return true;
  };

  const textContentClean = (el) => {
    // Collapse whitespace the way visual text does.
    return el.textContent.replace(/\s+/g, ' ').trim();
  };

  // Effective z-index: the nearest non-auto z-index on the element or
  // an ancestor below the section. PPTX stacking is purely shape
  // order, so primitives are stable-sorted by this value after the
  // walk — that's what keeps a z-index:2 label painting above a
  // sibling card that comes later in document order.
  const effectiveZ = (el, sectionEl) => {
    let cur = el;
    while (cur && cur !== sectionEl) {
      const z = window.getComputedStyle(cur).zIndex;
      if (z !== 'auto') {
        const n = parseInt(z, 10);
        if (!isNaN(n)) return n;
      }
      cur = cur.parentElement;
    }
    return 0;
  };

  // Rotation (deg, clockwise-positive like CSS) from a computed
  // transform matrix. Pure translations return 0.
  const getRotationDeg = (tr) => {
    if (!tr || tr === 'none') return 0;
    const m = tr.match(/^matrix\(([^)]+)\)/);
    if (!m) return 0;   // matrix3d etc — not handled
    const parts = m[1].split(',').map(parseFloat);
    const deg = Math.atan2(parts[1], parts[0]) * 180 / Math.PI;
    return Math.abs(deg) < 0.5 ? 0 : deg;
  };

  // Overflow clipping. HTML cards clip their children (colored header
  // cells, accent bars) to a rounded outline via overflow:hidden +
  // border-radius. PPTX has no shape clipping, so square-cornered
  // child rects poke past the card silhouette. We track the
  // intersection of clipping-ancestor boxes down the walk, intersect
  // every rect with it, and round/inset corners where a rect hugs a
  // rounded clip edge.
  const OVERFLOW_CLIPS = new Set(['hidden', 'clip', 'auto', 'scroll']);
  const intersectClip = (px, py, pw, ph, clip) => {
    if (!clip) return { x: px, y: py, w: pw, h: ph };
    const x0 = Math.max(px, clip.x0), y0 = Math.max(py, clip.y0);
    const x1 = Math.min(px + pw, clip.x1), y1 = Math.min(py + ph, clip.y1);
    if (x1 - x0 < 0.5 || y1 - y0 < 0.5) return null;
    return { x: x0, y: y0, w: x1 - x0, h: y1 - y0 };
  };
  // Per-corner border radii. Same-side pairs map to native pptx
  // shapes (round2SameRect); a bar authored "radius: 20 0 0 20" to
  // follow its card's left curve must not export as an all-corner
  // pill.
  const readRadii = (cs) => {
    const tl = parseFloat(cs.borderTopLeftRadius) || 0;
    const tr = parseFloat(cs.borderTopRightRadius) || 0;
    const br = parseFloat(cs.borderBottomRightRadius) || 0;
    const bl = parseFloat(cs.borderBottomLeftRadius) || 0;
    return { tl, tr, br, bl, max: Math.max(tl, tr, br, bl) };
  };

  // Apply clip + rounded-corner adaptation to a rect prim, then push.
  const pushClippedRect = (primitives, prim, clip) => {
    const r = intersectClip(prim.x, prim.y, prim.w, prim.h, clip);
    if (!r) return;
    prim.x = r.x; prim.y = r.y; prim.w = r.w; prim.h = r.h;
    const r4 = prim.radii;
    delete prim.radii;
    let adapted = false;
    if (clip && clip.radius >= 3) {
      const t = 3.5;   // borders make children sit 1-2px inside the clip box
      const hugsTop = Math.abs(r.y - clip.y0) < t;
      const hugsBottom = Math.abs(r.y + r.h - clip.y1) < t;
      const hugsLeft = Math.abs(r.x - clip.x0) < t;
      const hugsRight = Math.abs(r.x + r.w - clip.x1) < t;
      const spansW = hugsLeft && hugsRight;
      const spansH = hugsTop && hugsBottom;
      if (spansH && (hugsLeft || hugsRight) && r.w <= clip.radius) {
        // Thin vertical accent bar at a rounded edge (any authored
        // radius) — pull it in from the curved corners instead of
        // poking past them, and square it off.
        prim.y += clip.radius; prim.h -= 2 * clip.radius;
        prim.radius = 0;
        if (prim.h < 0.5) return;
        adapted = true;
      } else if (spansW && (hugsTop || hugsBottom) && r.h <= clip.radius) {
        // Thin horizontal strip at a rounded edge — same treatment.
        prim.x += clip.radius; prim.w -= 2 * clip.radius;
        prim.radius = 0;
        if (prim.w < 0.5) return;
        adapted = true;
      } else if (!prim.radius || prim.radius < 3) {
        if (spansW && spansH) {
          prim.radius = clip.radius;          // covers the whole card
          adapted = true;
        } else if (spansW && hugsTop) {
          prim.corner = 'top'; prim.radius = clip.radius;
          adapted = true;
        } else if (spansW && hugsBottom) {
          prim.corner = 'bottom'; prim.radius = clip.radius;
          adapted = true;
        }
      }
    }
    // Authored same-side corner pairs (e.g. radius: 20 20 0 0 on a
    // header band) → native top/bottom-rounded shape instead of an
    // all-corner rounded rect.
    if (!adapted && r4 && r4.max >= 3) {
      const on = (v) => v >= 3;
      if (on(r4.tl) && on(r4.tr) && !on(r4.br) && !on(r4.bl)) {
        prim.corner = 'top'; prim.radius = Math.max(r4.tl, r4.tr);
      } else if (!on(r4.tl) && !on(r4.tr) && on(r4.br) && on(r4.bl)) {
        prim.corner = 'bottom'; prim.radius = Math.max(r4.bl, r4.br);
      } else if (on(r4.tl) && on(r4.bl) && !on(r4.tr) && !on(r4.br)) {
        // Left pair (radius: N 0 0 N) — first segment of a stacked
        // bar. Uniform rounding notched the flush seam against the
        // next segment.
        prim.corner = 'left'; prim.radius = Math.max(r4.tl, r4.bl);
      } else if (!on(r4.tl) && !on(r4.bl) && on(r4.tr) && on(r4.br)) {
        prim.corner = 'right'; prim.radius = Math.max(r4.tr, r4.br);
      } else if (!(on(r4.tl) && on(r4.tr) && on(r4.br) && on(r4.bl))) {
        // Mixed singles (one rounded corner etc.) — uniform rounding
        // would notch flush seams; squaring off is the lesser error.
        prim.radius = 0;
      }
    }
    primitives.push(prim);
  };

  // Parse the first gradient layer out of a computed background-image.
  // Computed form is normalized by the browser, e.g.
  //   linear-gradient(rgb(0, 0, 0) 0%, rgb(255, 255, 255) 100%)
  //   repeating-linear-gradient(-45deg, rgb(229, 72, 61) 0px, ... 18px)
  //   radial-gradient(at 80% 20%, rgba(0, 113, 227, 0.25) 0%, transparent 60%)
  // Returns { kind: 'pattern'|'linear'|'radial', ... } or null.
  const splitTopLevel = (s, sep) => {
    const out = [];
    let depth = 0, cur = '';
    for (const ch of s) {
      if (ch === '(') depth++;
      if (ch === ')') depth--;
      if (ch === sep && depth === 0) { out.push(cur); cur = ''; continue; }
      cur += ch;
    }
    if (cur.trim()) out.push(cur);
    return out;
  };
  const parseBgGradient = (bgImage) => {
    if (!bgImage || bgImage === 'none') return null;
    const layer = splitTopLevel(bgImage, ',').find((l) => l.includes('gradient('));
    if (!layer) return null;
    const m = layer.match(/(repeating-linear|linear|radial)-gradient\((.*)\)\s*$/s);
    if (!m) return null;
    const kind = m[1];
    const args = splitTopLevel(m[2], ',').map((x) => x.trim());
    let angleDeg = 180;   // CSS default: to bottom
    let atPos = null;
    let stopArgs = args;
    if (args.length && /^-?[\d.]+deg$/.test(args[0])) {
      angleDeg = parseFloat(args[0]);
      stopArgs = args.slice(1);
    } else if (args.length && args[0].startsWith('to ')) {
      const dir = args[0].slice(3).trim();
      angleDeg = { 'top': 0, 'right': 90, 'bottom': 180, 'left': 270,
                   'top right': 45, 'right top': 45, 'bottom right': 135, 'right bottom': 135,
                   'bottom left': 225, 'left bottom': 225, 'top left': 315, 'left top': 315 }[dir] ?? 180;
      stopArgs = args.slice(1);
    } else if (args.length && (args[0].includes('at ') || /^(circle|ellipse|closest|farthest)/.test(args[0]))) {
      const am = args[0].match(/at\s+([\d.]+)%\s+([\d.]+)%/);
      if (am) atPos = { x: parseFloat(am[1]) / 100, y: parseFloat(am[2]) / 100 };
      stopArgs = args.slice(1);
    }
    const stops = [];
    for (const sa of stopArgs) {
      let col = null, alpha = 1;
      const cm = sa.match(/rgba?\([^)]+\)/);
      if (cm) {
        const pc = parseColor(cm[0]);
        if (pc) { col = toHex(pc); alpha = pc.alpha; }
        else { col = null; alpha = 0; }   // rgba(...,0) → transparent stop
      } else if (/\btransparent\b/.test(sa)) {
        col = null; alpha = 0;
      } else {
        continue;
      }
      const pm = sa.match(/(-?[\d.]+)(%|px)\s*$/);
      const pos = pm ? { v: parseFloat(pm[1]), unit: pm[2] } : null;
      stops.push({ color: col, alpha, pos });
    }
    if (stops.length < 2) return null;
    // Normalize positions to 0..1 (px positions normalize against the
    // largest px seen — good enough for stripe patterns and bars).
    const maxPx = Math.max(1, ...stops.map((s) => (s.pos && s.pos.unit === 'px') ? s.pos.v : 0));
    stops.forEach((s, i) => {
      if (!s.pos) s.f = i / (stops.length - 1);
      else s.f = s.pos.unit === '%' ? s.pos.v / 100 : s.pos.v / maxPx;
      delete s.pos;
    });
    if (kind === 'repeating-linear') {
      // Hatch stripes → pptx pattern fill. Use the two most distinct colors.
      const opaque = stops.filter((s) => s.color);
      if (!opaque.length) return null;
      const first = opaque[0];
      const other = opaque.find((s) => s.color !== first.color);
      // No second opaque color (e.g. translucent-white stripes over
      // transparent) → bg resolves to the underlying base color later.
      return { kind: 'pattern', angle: ((angleDeg % 180) + 180) % 180,
               fg: first.color, fgAlpha: first.alpha,
               bg: other ? other.color : null,
               bgAlpha: other ? other.alpha : 0 };
    }
    return { kind: kind === 'radial' ? 'radial' : 'linear', angle: angleDeg, at: atPos, stops };
  };

  const hexToRgbObj = (hex) => ({
    r: parseInt(hex.slice(0, 2), 16), g: parseInt(hex.slice(2, 4), 16),
    b: parseInt(hex.slice(4, 6), 16), alpha: 1,
  });

  // Composite a gradient's translucent stops over the layer beneath it
  // (the element's own background-color, else the nearest opaque
  // ancestor). A texture like rgba(255,255,255,.18) stripes over a red
  // bar must export as light-red-on-red — exporting the raw stripe
  // colors loses the bar entirely.
  const resolveGradOverBase = (grad, baseRgb) => {
    if (!grad) return null;
    const blendHex = (hex, alpha) => {
      if (!hex) return toHex(baseRgb);
      if (alpha >= 0.999) return hex;
      return blendOverBase({ ...hexToRgbObj(hex), alpha }, baseRgb);
    };
    if (grad.kind === 'pattern') {
      return { kind: 'pattern', angle: grad.angle,
               fg: blendHex(grad.fg, grad.fgAlpha ?? 1),
               bg: blendHex(grad.bg, grad.bgAlpha ?? 1) };
    }
    return { kind: grad.kind, angle: grad.angle, at: grad.at,
             stops: grad.stops.map((s) => ({ color: blendHex(s.color, s.alpha), alpha: 1, f: s.f })) };
  };

  // Rasterize an inline <svg> to a PNG data URL via canvas. Keeps
  // icons visible in the export (they'd otherwise vanish — svg has no
  // text/img/bg for the walker to capture). Computed presentation
  // styles are inlined onto a clone first so stylesheet-applied
  // stroke-width and currentColor survive standalone rendering.
  const svgCache = new Map();
  const rasterizeSvg = async (el, w, h) => {
    const color = window.getComputedStyle(el).color;
    const key = el.outerHTML + '|' + color + '|' + Math.round(w) + 'x' + Math.round(h);
    if (svgCache.has(key)) return svgCache.get(key);
    let url = null;
    try {
      const clone = el.cloneNode(true);
      const srcEls = [el, ...el.querySelectorAll('*')];
      const dstEls = [clone, ...clone.querySelectorAll('*')];
      for (let i = 0; i < srcEls.length; i++) {
        const cs = window.getComputedStyle(srcEls[i]);
        const d = dstEls[i];
        if (!d.setAttribute) continue;
        d.setAttribute('stroke', cs.stroke);
        d.setAttribute('fill', cs.fill);
        d.setAttribute('stroke-width', cs.strokeWidth);
        d.setAttribute('stroke-linecap', cs.strokeLinecap);
        d.setAttribute('stroke-linejoin', cs.strokeLinejoin);
        if (cs.strokeDasharray && cs.strokeDasharray !== 'none') d.setAttribute('stroke-dasharray', cs.strokeDasharray);
        if (cs.strokeDashoffset && cs.strokeDashoffset !== '0px') d.setAttribute('stroke-dashoffset', cs.strokeDashoffset);
        if (cs.opacity && cs.opacity !== '1') d.setAttribute('opacity', cs.opacity);
        if (i > 0 && cs.fontSize) d.setAttribute('font-size', cs.fontSize);
        if (i > 0 && cs.fontFamily) d.setAttribute('font-family', cs.fontFamily);
        if (i > 0 && cs.fontWeight) d.setAttribute('font-weight', cs.fontWeight);
      }
      clone.setAttribute('xmlns', 'http://www.w3.org/2000/svg');
      clone.setAttribute('width', String(Math.max(1, Math.round(w))));
      clone.setAttribute('height', String(Math.max(1, Math.round(h))));
      const xml = new XMLSerializer().serializeToString(clone);
      const svgUri = 'data:image/svg+xml;base64,' + btoa(unescape(encodeURIComponent(xml)));
      const img = new Image();
      img.src = svgUri;
      await img.decode();
      const scale = 3;   // rasterize at 3x for crispness when resized
      const canvas = document.createElement('canvas');
      canvas.width = Math.max(1, Math.round(w * scale));
      canvas.height = Math.max(1, Math.round(h * scale));
      canvas.getContext('2d').drawImage(img, 0, 0, canvas.width, canvas.height);
      url = canvas.toDataURL('image/png');
    } catch (e) {
      url = null;
    }
    svgCache.set(key, url);
    return url;
  };

  // Capture ::before / ::after pseudo-elements. The browser already
  // laid these out as part of the parent's render, but they aren't in
  // the DOM tree — so walkSlide() would skip them entirely. Many
  // systems use pseudos for accent rules (.accent-rule::before),
  // bullet markers (li::before), drop caps, etc. — losing them tanks
  // visual fidelity.
  //
  // We handle the common case: position: absolute relative to the
  // parent (which most decorative pseudos use). Counter-based content
  // (ol.numbered li::before { content: counter(item) }) is captured by
  // manually computing the counter value as we traverse.
  const inspectPseudo = (el, which) => {
    const cs = window.getComputedStyle(el, which);
    if (cs.display === 'none' || cs.visibility === 'hidden' || parseFloat(cs.opacity) === 0) return null;
    let content = cs.content;
    if (!content || content === 'none' || content === 'normal') return null;

    // Resolve the content string. Three cases worth handling:
    //   "string"    → quoted literal, strip quotes
    //   counter(x)  → use stored counter value (computed below)
    //   anything else (image, attr()) → skip
    let displayText = '';
    if (content.startsWith('"') || content.startsWith("'")) {
      displayText = content.slice(1, -1);
    } else if (/^counter/i.test(content)) {
      // We can't evaluate counter() from JS — but if the parent is
      // inside an ol, we can read the parent's index manually. Best
      // effort: try the closest ol ancestor.
      const ol = el.closest('ol');
      if (ol) {
        const items = Array.from(ol.querySelectorAll(':scope > li'));
        const idx = items.indexOf(el) + 1;
        if (idx > 0) {
          // Detect style — decimal vs decimal-leading-zero vs upper-roman.
          const m = content.match(/counter\([^,)]+(?:,\s*(\S+?))?\s*\)/i);
          const style = m && m[1] ? m[1].toLowerCase() : 'decimal';
          if (style.includes('roman')) {
            displayText = toRoman(idx, style.includes('upper'));
          } else if (style.includes('leading-zero')) {
            displayText = String(idx).padStart(2, '0');
          } else {
            displayText = String(idx);
          }
        }
      }
      if (!displayText) return null;
    } else if (content === '""' || content === "''") {
      displayText = '';
    } else {
      return null;
    }

    const parentRect = el.getBoundingClientRect();
    const pos = cs.position;
    const top = parseFloat(cs.top);
    const left = parseFloat(cs.left);
    const right = parseFloat(cs.right);
    const bottom = parseFloat(cs.bottom);
    const width = parseFloat(cs.width);
    const height = parseFloat(cs.height);
    const bgColor = parseColor(cs.backgroundColor);
    const fontSizePx = parseFloat(cs.fontSize) || 0;
    const colorRaw = parseColor(cs.color);

    // We handle position: absolute relative to the parent (which is
    // typically position: relative). Inline pseudos that flow with
    // text need browser layout info we can't easily get from JS — skip
    // for now. Static block pseudos (position: static) we approximate
    // by treating the parent's top-left as the anchor.
    let px, py, pw, ph;
    if (pos === 'absolute' || pos === 'static' || pos === 'relative') {
      const safeTop  = isNaN(top)  ? 0 : top;
      const safeLeft = isNaN(left) ? 0 : left;
      px = parentRect.left + safeLeft;
      py = parentRect.top  + safeTop;
      if (!isNaN(width)) {
        pw = width;
      } else if (!isNaN(left) && !isNaN(right)) {
        pw = parentRect.width - left - right;
      } else if (displayText) {
        // Text content — width is roughly textLen * fontSize * 0.55
        pw = displayText.length * fontSizePx * 0.6 + 4;
      } else {
        pw = 0;
      }
      if (!isNaN(height)) {
        ph = height;
      } else if (displayText) {
        ph = fontSizePx * 1.4;
      } else {
        ph = 0;
      }
    } else {
      return null;
    }
    if (pw < 0.5 || ph < 0.5) return null;
    return {
      x: px, y: py, w: pw, h: ph,
      bg: bgColor,
      color: colorRaw,
      text: displayText,
      fontSizePx,
      fontFamily: cs.fontFamily,
      fontWeight: parseInt(cs.fontWeight) || 400,
      fontStyle: cs.fontStyle,
      letterSpacing: cs.letterSpacing === 'normal' ? 0 : parseFloat(cs.letterSpacing) || 0,
      textTransform: cs.textTransform || 'none',
      radius: readRadii(cs).max,
      radii: readRadii(cs),
    };
  };

  // Tiny roman-numeral helper for counter() pseudos that use upper-roman.
  const toRoman = (n, upper) => {
    const m = [['M',1000],['CM',900],['D',500],['CD',400],['C',100],['XC',90],['L',50],['XL',40],['X',10],['IX',9],['V',5],['IV',4],['I',1]];
    let s = '';
    for (const [r, v] of m) { while (n >= v) { s += r; n -= v; } }
    return upper === false ? s.toLowerCase() : s;
  };

  // Per-segment style read from the text node's nearest element
  // ancestor. Lets `<p>foo <strong>bar</strong></p>` emit two runs in
  // the pptx paragraph — one regular, one bold — instead of one run
  // with the paragraph's base weight.
  const CSS_FONT_KEYWORDS = new Set([
    'sans-serif', 'serif', 'monospace', 'cursive', 'fantasy',
    'system-ui', 'ui-sans-serif', 'ui-serif', 'ui-monospace', 'ui-rounded',
    'inherit', 'initial', 'unset',
  ]);
  // Pick the first *actually available* named font in the stack —
  // document.fonts.check() resolves both locally-installed fonts and
  // loaded webfonts. Without the check, a stack like
  // "SF Mono", Menlo, monospace exports as "SF Mono" even on machines
  // that don't have it, and PowerPoint/Keynote flag a missing font.
  const fontPickCache = new Map();
  const fontAvailable = (name) => {
    try { return document.fonts.check(`12px "${name}"`); } catch (e) { return true; }
  };
  const pickFontFamily = (s) => {
    if (fontPickCache.has(s)) return fontPickCache.get(s);
    const named = (s || '').split(',').map(x => x.replace(/['"]/g, '').trim())
      .filter(x => x && !CSS_FONT_KEYWORDS.has(x.toLowerCase()) && !x.startsWith('-'));
    // Apple's SF-prefixed system fonts (SF Mono, SF Pro Display, ...)
    // pass document.fonts.check in chromium but are NOT selectable in
    // Keynote/PowerPoint — exporting them triggers missing-font
    // fallback in the very apps this file targets. Prefer the first
    // office-usable name in the stack (e.g. Menlo over SF Mono).
    const officeSafe = named.filter(x => !/^SF[ -]/i.test(x) && !x.startsWith('.'));
    const picked = officeSafe.find(fontAvailable) || named.find(fontAvailable)
      || officeSafe[0] || named[0] || 'Helvetica Neue';
    fontPickCache.set(s, picked);
    return picked;
  };
  const styleOfTextNode = (n) => {
    // The text node's immediate element parent owns its styling.
    const parent = n.parentElement;
    if (!parent) return null;
    const cs = window.getComputedStyle(parent);
    const colorRaw = parseColor(cs.color);
    const colorBase = colorRaw ? effectiveBg(parent) : null;
    const tdLine = (cs.textDecorationLine || cs.textDecoration || '').toLowerCase();
    return {
      fontFamily:   pickFontFamily(cs.fontFamily),
      fontSizePx:   parseFloat(cs.fontSize) || 14,
      fontWeight:   parseInt(cs.fontWeight) || 400,
      fontStyle:    cs.fontStyle || 'normal',
      color:        colorRaw ? blendOverBase(colorRaw, colorBase) : '000000',
      letterSpacingPx: cs.letterSpacing === 'normal' ? 0 : (parseFloat(cs.letterSpacing) || 0),
      underline:    tdLine.includes('underline'),
      strikethrough: tdLine.includes('line-through'),
      textTransform: cs.textTransform || 'none',
    };
  };
  const styleKey = (s) =>
    !s ? '' : [s.fontFamily, s.fontSizePx, s.fontWeight, s.fontStyle, s.color,
               s.letterSpacingPx, s.underline, s.strikethrough, s.textTransform].join('|');
  const applyXform = (s, text) =>
      s === 'uppercase' ? text.toUpperCase()
    : s === 'lowercase' ? text.toLowerCase()
    : s === 'capitalize' ? text.replace(/\b\w/g, c => c.toUpperCase())
    : text;

  // Walk a leaf text element and return one shape-per-visual-line,
  // where each line is split into STYLED SEGMENTS at inline-element
  // boundaries. The browser already laid out the line; we just stitch
  // word-rects from text-node ranges back into segments per parent.
  //
  // The result lets PPTX emit one textbox per line, containing one
  // paragraph with N runs (one per style transition). Bold/italic/
  // colored spans inside a paragraph survive the trip.
  //
  // Whitespace correctness matters here: the word regex /\S+\s*/g
  // keeps TRAILING whitespace with each word but skips LEADING
  // whitespace — and inline-element boundaries routinely put the
  // separating space at the START of the following text node
  // ("<strong>bold</strong> rest"), or in a whitespace-only node
  // between elements. Both used to vanish, gluing words together
  // ("bold" + "rest" → "boldrest"). pendingWs tracks skipped leading
  // whitespace and re-inserts a single space when the line continues.
  // directOnly limits capture to text nodes whose parent IS el — used
  // for containers with mixed content (e.g. an inline-flex pill whose
  // children blockify: <span class="pill"><span class="dot"/>label</span>).
  // Element children get their own recursion; the bare label text node
  // belongs to the container itself and must be captured here or lost.
  const getTextLines = (el, directOnly = false) => {
    const walker = document.createTreeWalker(el, NodeFilter.SHOW_TEXT, null);
    const lines = [];   // { x, y, w, h, segments: [{text, style}] }
    let curLine = null;
    let curSeg = null;
    let curStyleKey = null;
    let pendingWs = false;
    let n;
    const flushSeg = () => {
      if (curSeg && curSeg.text && curLine) {
        curSeg.text = applyXform(curSeg.style.textTransform, curSeg.text);
        curLine.segments.push(curSeg);
      }
      curSeg = null; curStyleKey = null;
    };
    const flushLine = () => {
      flushSeg();
      if (curLine && curLine.segments.length) lines.push(curLine);
      curLine = null;
    };
    while ((n = walker.nextNode())) {
      if (directOnly && n.parentElement !== el) continue;
      const text = n.nodeValue;
      if (!text) continue;
      if (!text.trim()) {
        // Whitespace-only node between inline elements — remember it
        // so the next word on the same line gets its separator back.
        if (curLine) pendingWs = true;
        continue;
      }
      const segStyle = styleOfTextNode(n);
      const segKey = styleKey(segStyle);
      const re = /\S+\s*/g;
      let m;
      let firstMatchInNode = true;
      while ((m = re.exec(text)) !== null) {
        // Leading whitespace at the head of the node was skipped by
        // the regex; only the first match can have skipped any (later
        // gaps ride along as the previous match's trailing \s*).
        if (firstMatchInNode && m.index > 0) pendingWs = true;
        firstMatchInNode = false;
        const r = document.createRange();
        r.setStart(n, m.index);
        r.setEnd(n, m.index + m[0].length);
        const rects = r.getClientRects();
        for (const rect of rects) {
          if (rect.width === 0 || rect.height === 0) continue;
          // Same line if y is within a few px of the running line.
          const sameLine = curLine && Math.abs(rect.top - curLine.y) < Math.max(2, curLine.h * 0.4);
          if (!sameLine) {
            flushLine();
            curLine = { x: rect.left, y: rect.top, w: rect.width, h: rect.height, segments: [] };
            curSeg = { text: m[0], style: segStyle };
            curStyleKey = segKey;
          } else {
            curLine.w = Math.max(curLine.w, rect.right - curLine.x);
            curLine.h = Math.max(curLine.h, rect.height);
            const sep = pendingWs ? ' ' : '';
            if (segKey !== curStyleKey) {
              flushSeg();
              curSeg = { text: sep + m[0], style: segStyle };
              curStyleKey = segKey;
            } else {
              curSeg.text += sep + m[0];
            }
          }
          pendingWs = false;
          break;
        }
      }
    }
    flushLine();
    // Trim trailing whitespace from the last segment of each line.
    for (const ln of lines) {
      if (ln.segments.length) {
        const last = ln.segments[ln.segments.length - 1];
        last.text = last.text.replace(/\s+$/, '');
      }
      ln.segments = ln.segments.filter(s => s.text);
    }
    return lines.filter(l => l.segments.length);
  };

  const walkSlide = async (section) => {
    const primitives = [];
    const sectionRect = section.getBoundingClientRect();
    const sx = sectionRect.left, sy = sectionRect.top;

    // Canvas-measured font metrics, used to place text boxes by
    // BASELINE rather than line-box top. Browser line rects are
    // line-height tall (glyphs centered via half-leading); pptx draws
    // the first baseline at boxTop + font ascent. The mismatch grows
    // with font size — at display sizes (150px+ numerals with
    // line-height 0.9) lines land visibly wrong without this.
    const measureCtx = document.createElement('canvas').getContext('2d');
    const metricsCache = new Map();
    const fontMetrics = (family, sizePx, weight, style) => {
      const key = family + '|' + sizePx + '|' + weight + '|' + style;
      if (metricsCache.has(key)) return metricsCache.get(key);
      let out;
      try {
        measureCtx.font = `${style === 'italic' ? 'italic ' : ''}${weight} ${sizePx}px "${family}"`;
        const m = measureCtx.measureText('Hg');
        out = {
          asc: m.fontBoundingBoxAscent ?? sizePx * 0.8,
          desc: m.fontBoundingBoxDescent ?? sizePx * 0.2,
        };
      } catch (e) {
        out = { asc: sizePx * 0.8, desc: sizePx * 0.2 };
      }
      metricsCache.set(key, out);
      return out;
    };

    const segmentToRun = (seg) => ({
      text: seg.text,
      font_family:       seg.style.fontFamily,
      font_size_px:      seg.style.fontSizePx,
      font_weight:       seg.style.fontWeight,
      font_style:        seg.style.fontStyle,
      color:             seg.style.color,
      letter_spacing_px: seg.style.letterSpacingPx,
      underline:         seg.style.underline,
      strikethrough:     seg.style.strikethrough,
    });

    // One textbox per captured line. Modest width slack absorbs
    // renderer-level metric differences (the page is re-measured in
    // the export font before walking, so font-substitution drift is
    // already gone); the box anchors on the side the text aligns to
    // so slack grows away from the visual anchor.
    const emitTextLines = (lines, align, z) => {
      for (const ln of lines) {
        const slack = ln.w * 0.06 + 6;
        let bx = ln.x - sx;
        let balign = 'left';
        if (align === 'right') { bx -= slack; balign = 'right'; }
        else if (align === 'center') { bx -= slack / 2; balign = 'center'; }
        // Baseline correction: the browser centers glyphs in the
        // line-height box (half-leading); pptx puts the first baseline
        // at boxTop + ascent. Setting boxTop = lineTop + halfLeading
        // makes both baselines coincide (same font both sides, since
        // the page re-measures in the export font).
        let maxAsc = 0, maxDesc = 0;
        for (const seg of ln.segments) {
          const fm = fontMetrics(seg.style.fontFamily, seg.style.fontSizePx,
                                 seg.style.fontWeight, seg.style.fontStyle);
          maxAsc = Math.max(maxAsc, fm.asc);
          maxDesc = Math.max(maxDesc, fm.desc);
        }
        const halfLeading = (ln.h - (maxAsc + maxDesc)) / 2;
        primitives.push({
          kind: 'text',
          x: bx, y: ln.y - sy + halfLeading,
          w: ln.w + slack,
          h: maxAsc + maxDesc + 2,
          align: balign,
          z,
          runs: ln.segments.map(segmentToRun),
        });
      }
    };

    const recurse = async (el, depth, clip) => {
      const st = window.getComputedStyle(el);
      // display:contents generates no box (rect is 0×0), so the
      // isVisible gate below would kill the whole subtree — but its
      // children are real boxes hoisted into the parent. Walk through
      // it transparently.
      if (st.display === 'contents') {
        for (const child of el.children) await recurse(child, depth, clip);
        const hasDirectText = Array.from(el.childNodes).some(
          (c) => c.nodeType === Node.TEXT_NODE && c.nodeValue && c.nodeValue.trim());
        if (hasDirectText) {
          const align = st.textAlign === 'start' ? 'left' : st.textAlign;
          emitTextLines(getTextLines(el, true), align, effectiveZ(el, section));
        }
        return;
      }
      if (!isVisible(el, st)) return;
      const r = el.getBoundingClientRect();
      const x = r.left - sx, y = r.top - sy, w = r.width, h = r.height;
      const z = effectiveZ(el, section);

      // Inline SVG (icon) → rasterize to a small PNG picture. SVGs
      // have no text nodes / img / bg the walker could capture; they
      // used to vanish entirely (empty glyph boxes everywhere).
      if (el.tagName && el.tagName.toLowerCase() === 'svg') {
        const dataUrl = await rasterizeSvg(el, w, h);
        if (dataUrl) primitives.push({ kind: 'pic', x, y, w, h, src: dataUrl, z });
        return;
      }

      // Image element → picture primitive.
      if (el.tagName === 'IMG' && (el.currentSrc || el.src)) {
        primitives.push({ kind: 'pic', x, y, w, h, src: el.currentSrc || el.src, z });
        return;
      }

      // Background fill or border on a container → rect primitive(s).
      const bg = parseColor(st.backgroundColor);
      const sides = [
        { side: 'top',    w: parseFloat(st.borderTopWidth)    || 0, c: parseColor(st.borderTopColor) },
        { side: 'right',  w: parseFloat(st.borderRightWidth)  || 0, c: parseColor(st.borderRightColor) },
        { side: 'bottom', w: parseFloat(st.borderBottomWidth) || 0, c: parseColor(st.borderBottomColor) },
        { side: 'left',   w: parseFloat(st.borderLeftWidth)   || 0, c: parseColor(st.borderLeftColor) },
      ];
      const visibleSides = sides.filter(s => s.w > 0.5 && s.c);
      const totalBw = sides.reduce((a, s) => a + s.w, 0);
      const radii = readRadii(st);
      const radius = radii.max;
      // Gradient / hatch-pattern backgrounds (progress bars, discard
      // segments, glow washes) come from background-image, not
      // background-color — they used to emit nothing and leave holes.
      const grad = parseBgGradient(st.backgroundImage);

      // This element's children get clipped here if it sets overflow.
      // The element's OWN bg/border rect is clipped by the inherited
      // clip only (a box never clips itself).
      let childClip = clip;
      if (OVERFLOW_CLIPS.has(st.overflowX) || OVERFLOW_CLIPS.has(st.overflowY)) {
        const nc = { x0: x, y0: y, x1: x + w, y1: y + h, radius };
        childClip = clip
          ? { x0: Math.max(clip.x0, nc.x0), y0: Math.max(clip.y0, nc.y0),
              x1: Math.min(clip.x1, nc.x1), y1: Math.min(clip.y1, nc.y1), radius: nc.radius }
          : nc;
      }

      if (bg || grad || totalBw > 0.5) {
        const base = effectiveBg(el);
        // Gradient layers paint over the element's own bg color (if
        // any), which itself paints over the ancestor base.
        const gradBase = bg ? hexToRgbObj(blendOverBase(bg, base)) : base;
        const gradResolved = resolveGradOverBase(grad, gradBase);
        // Detect uniform border (all 4 sides identical width + color).
        // PPTX shape.line handles this natively; asymmetric borders
        // need separate overlay rects per side.
        const sideKey = (s) => s.c
          ? `${s.w}|${s.c.r}|${s.c.g}|${s.c.b}|${s.c.alpha.toFixed(2)}`
          : `${s.w}|null`;
        const allUniform = visibleSides.length === 4 &&
          visibleSides.every(s => sideKey(s) === sideKey(visibleSides[0]));

        if (allUniform || visibleSides.length === 0 || !bg && visibleSides.length === 1) {
          // Single rect with built-in line. Either uniform border, or
          // no border, or a single-side border on a no-fill container
          // (rare — we'd lose this case; fall through to overlay below).
          const ref = visibleSides[0];
          pushClippedRect(primitives, {
            kind: 'rect', x, y, w, h,
            fill:   bg ? blendOverBase(bg, base) : null,
            grad: gradResolved,
            line:   (allUniform && ref) ? blendOverBase(ref.c, base) : null,
            line_w: (allUniform && ref) ? ref.w : 0,
            radius, radii,
            z,
          }, clip);
        }
        if (!allUniform && visibleSides.length > 0) {
          // Asymmetric border: emit the fill rect (no line) plus
          // per-side overlay rects. This handles patterns like
          // border-left: 4px solid var(--accent) used on alerts,
          // accent cards, and the editorial drop-cap bar.
          if (allUniform === false && (bg || grad)) {
            pushClippedRect(primitives, {
              kind: 'rect', x, y, w, h,
              fill: bg ? blendOverBase(bg, base) : null, grad: gradResolved,
              line: null, line_w: 0, radius, radii,
              z,
            }, clip);
          } else if (allUniform === false && !bg) {
            // Empty rect just for the border? No — skip the fill rect
            // and only emit the side overlays.
          }
          for (const s of visibleSides) {
            const stroke = blendOverBase(s.c, base);
            let bx, by, bw, bh;
            if (s.side === 'top')    { bx = x;          by = y;          bw = w;   bh = s.w; }
            if (s.side === 'right')  { bx = x + w - s.w; by = y;          bw = s.w; bh = h;  }
            if (s.side === 'bottom') { bx = x;          by = y + h - s.w; bw = w;   bh = s.w; }
            if (s.side === 'left')   { bx = x;          by = y;          bw = s.w; bh = h;  }
            // A border strip on a rounded element pokes past the
            // curved corners (the real border follows the curve).
            // Inset the strip along its length by the radius.
            if (radius >= 3) {
              if (s.side === 'top' || s.side === 'bottom') { bx += radius; bw -= 2 * radius; }
              else { by += radius; bh -= 2 * radius; }
              if (bw < 0.5 || bh < 0.5) continue;
            }
            pushClippedRect(primitives, {
              kind: 'rect', x: bx, y: by, w: bw, h: bh,
              fill: stroke, line: null, line_w: 0, radius: 0,
              z,
            }, clip);
          }
        }
      }

      // Pseudo-elements (::before / ::after). Inspect BEFORE the
      // text-leaf early return below — list markers (ol.numbered
      // li::before counters, ul.bullets li::before dashes) hang off
      // elements whose DOM content is pure text, and used to be
      // skipped entirely.
      for (const which of ['::before', '::after']) {
        const ps = inspectPseudo(el, which);
        if (!ps) continue;
        const px = ps.x - sx, py = ps.y - sy;
        if (ps.bg) {
          const base = effectiveBg(el);
          // Pseudo content is clipped by the element's own overflow
          // (accent bars on overflow:hidden cards).
          pushClippedRect(primitives, {
            kind: 'rect', x: px, y: py, w: ps.w, h: ps.h,
            fill: blendOverBase(ps.bg, base),
            line: null, line_w: 0,
            radius: ps.radius || 0, radii: ps.radii,
            z,
          }, childClip);
        }
        if (ps.text) {
          const fontFamily = pickFontFamily(ps.fontFamily);
          const base = effectiveBg(el);
          const colorHex = ps.color ? blendOverBase(ps.color, base) : '000000';
          primitives.push({
            kind: 'text', x: px, y: py, w: ps.w + 4, h: ps.h,
            align: 'left',
            z,
            runs: [{
              text: applyXform(ps.textTransform, ps.text),
              font_family: fontFamily,
              font_size_px: ps.fontSizePx,
              font_weight: ps.fontWeight,
              font_style: ps.fontStyle,
              color: colorHex,
              letter_spacing_px: ps.letterSpacing,
              underline: false,
              strikethrough: false,
            }],
          });
        }
      }

      // Leaf-ish element with text → emit one text primitive PER
      // VISUAL LINE. Capturing per-line rects (instead of one box for
      // the whole element) means pptx never has to wrap our text:
      // each line is its own absolutely-positioned shape, sized for
      // exactly its rendered HTML width + height. Font metric drift
      // between HTML and pptx can't push a wrap onto a new line or
      // crash into the row below.
      if (hasOnlyTextOrInline(el)) {
        const txt = textContentClean(el);
        if (txt) {
          const align = st.textAlign === 'start' ? 'left' : st.textAlign;

          // Vertical writing mode (writing-mode: vertical-rl/lr — axis
          // labels etc): per-line capture produces garbage (each word
          // becomes a narrow stacked box). Emit ONE textbox covering
          // the element's post-transform bounds and let PPTX rotate
          // the text natively via bodyPr@vert.
          const wm = st.writingMode || '';
          const rotation = getRotationDeg(st.transform);
          if (wm.startsWith('vertical') || wm.startsWith('sideways') || rotation !== 0) {
            const baseStyle = (() => {
              // Style from the first text node, falling back to the element.
              const tw = document.createTreeWalker(el, NodeFilter.SHOW_TEXT, null);
              let tn;
              while ((tn = tw.nextNode())) { if (tn.nodeValue && tn.nodeValue.trim()) return styleOfTextNode(tn); }
              return null;
            })();
            if (baseStyle) {
              const prim = {
                kind: 'text',
                align: 'center',
                anchor: 'middle',
                z,
                runs: [segmentToRun({ text: applyXform(baseStyle.textTransform, txt), style: baseStyle })],
              };
              if (wm.startsWith('vertical') || wm.startsWith('sideways')) {
                // vertical-rl reads top→bottom with glyphs rotated 90°
                // clockwise → pptx vert="vert". sideways-lr is the
                // 270° variant.
                prim.vert = (wm === 'sideways-lr') ? 'vert270' : 'vert';
                prim.x = x; prim.y = y; prim.w = w; prim.h = h;
              } else {
                // Pure CSS rotation: pptx rotates about the shape
                // center, same as CSS default transform-origin. Emit
                // the UNROTATED box centered where the rendered
                // (post-transform) box is centered.
                const ow = el.offsetWidth || w, oh = el.offsetHeight || h;
                prim.x = x + w / 2 - ow / 2;
                prim.y = y + h / 2 - oh / 2;
                prim.w = ow; prim.h = oh;
                prim.rot = rotation;
              }
              primitives.push(prim);
            }
            return;
          }

          // Per-line segmented capture. Each line becomes one textbox
          // primitive containing N styled runs (one per style change
          // boundary inside the line). Bold/italic/colored spans inside
          // a paragraph survive into pptx as separate runs sharing the
          // same paragraph + position.
          const lines = getTextLines(el);
          if (lines.length === 0) {
            // Fallback: emit a single run with the element's style.
            const baseStyle = styleOfTextNode({ parentElement: el }) || {
              fontFamily: pickFontFamily(st.fontFamily),
              fontSizePx: parseFloat(st.fontSize) || 14,
              fontWeight: parseInt(st.fontWeight) || 400,
              fontStyle: st.fontStyle || 'normal',
              color: '000000',
              letterSpacingPx: 0,
              underline: false, strikethrough: false,
              textTransform: st.textTransform || 'none',
            };
            primitives.push({
              kind: 'text', x, y, w: w * 1.05, h, align,
              z,
              runs: [segmentToRun({ text: applyXform(baseStyle.textTransform, txt), style: baseStyle })],
            });
          } else {
            emitTextLines(lines, align, z);
          }
        }
        return;
      }

      // Container — recurse into children, then pick up any text
      // nodes that are DIRECT children of this element. Mixed-content
      // containers (an inline-flex pill: <span class="pill"><span
      // class="dot"/>label</span>) blockify their element children, so
      // they fail hasOnlyTextOrInline — but the bare label text node
      // still belongs to the container and would otherwise vanish.
      for (const child of el.children) await recurse(child, depth + 1, childClip);
      const hasDirectText = Array.from(el.childNodes).some(
        (c) => c.nodeType === Node.TEXT_NODE && c.nodeValue && c.nodeValue.trim());
      if (hasDirectText) {
        const align = st.textAlign === 'start' ? 'left' : st.textAlign;
        emitTextLines(getTextLines(el, true), align, z);
      }
    };

    for (const child of section.children) await recurse(child, 0, null);

    // PPTX stacking is shape order. Stable-sort by effective z-index
    // so z-index:2 labels outpaint sibling cards that come later in
    // document order; equal z keeps document order (Array sort is
    // stable in modern V8).
    primitives.sort((a, b) => (a.z || 0) - (b.z || 0));

    return {
      width: sectionRect.width,
      height: sectionRect.height,
      bg: (() => {
        const bs = window.getComputedStyle(section);
        const c = parseColor(bs.backgroundColor);
        if (c && c.alpha >= 0.999) return toHex(c);
        // Section bg is transparent or translucent — fall through to body.
        const bb = parseColor(window.getComputedStyle(document.body).backgroundColor);
        if (c && bb) return blendOverBase(c, bb);
        return bb ? toHex(bb) : 'FFFFFF';
      })(),
      primitives,
    };
  };

  // Re-measure the page in the fonts the export will actually emit.
  // The browser typically renders -apple-system (SF Pro) while the
  // pptx says "Helvetica Neue" — a few percent of width drift per
  // line, which is enough to slide a line's tail underneath the next
  // shape (occlusion, not clipping: pptx z-order is paint order).
  // Overriding font-family up front makes the measured line rects,
  // wrap points, and flex/grid reflow all match the export font.
  document.querySelectorAll('deck-stage section, deck-stage section *').forEach((el) => {
    try {
      const fam = window.getComputedStyle(el).fontFamily;
      if (fam) el.style.fontFamily = '"' + pickFontFamily(fam) + '"';
    } catch (e) { /* non-styleable node */ }
  });
  // PPTX can't express OpenType features. Design systems set
  // font-feature-settings: "tnum" (tabular numerals), so the browser
  // would measure wide evenly-spaced digits while PowerPoint renders
  // proportional ones — digit-heavy display text ("10–20×", "132")
  // ends up with visibly different rhythm. Measure with features off
  // so geometry matches what the renderer will actually draw.
  const featKill = document.createElement('style');
  featKill.textContent = `deck-stage section, deck-stage section * {
    font-feature-settings: normal !important;
    font-variant-numeric: normal !important;
    font-variant-ligatures: normal !important;
  }`;
  document.head.appendChild(featKill);
  void document.body.offsetHeight;

  // Walk one section at a time: show it, then capture its primitives.
  // Showing one at a time keeps the original flex/grid layout intact
  // instead of all sections fighting for the same viewport.
  const result = [];
  for (let i = 0; i < sections.length; i++) {
    showOnly(i);
    result.push(await walkSlide(sections[i]));
  }
  return result;
})();
"""


def hex_to_rgb(hex_str: str) -> RGBColor:
    s = hex_str.lstrip('#').upper()
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def px_to_emu(px: float) -> int:
    return int(px * PX_TO_EMU)


def fetch_image(src: str) -> bytes | None:
    """Return image bytes for a src URL or data URI, or None on failure."""
    if src.startswith("data:"):
        # data:[<mime>][;base64],<data>
        head, _, payload = src.partition(",")
        if ";base64" in head:
            try: return base64.b64decode(payload)
            except Exception: return None
        # Non-base64 data URI — rarely used for images.
        return urllib.parse.unquote(payload).encode("latin-1", "ignore")
    try:
        with urllib.request.urlopen(src, timeout=15) as r:
            return r.read()
    except Exception:
        return None


def apply_gradient_fill(shp, grad: dict) -> None:
    """Replace the shape fill with a native pattern / gradient fill.

    pattern → a:pattFill with a hatch preset picked by stripe angle
    linear  → a:gradFill + a:lin (CSS angle is from-north, pptx from-east)
    radial  → a:gradFill + a:path[circle] with fillToRect at the CSS
              'at x% y%' position
    Transparent CSS stops carry over as srgbClr + a:alpha, which is the
    one place pptx transparency is actually honored everywhere.
    """
    sp_pr = shp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        for e in sp_pr.findall(qn(tag)):
            sp_pr.remove(e)

    def srgb_el(color: str | None, alpha: float) -> object:
        clr = sp_pr.makeelement(qn("a:srgbClr"), {"val": color or "FFFFFF"})
        if alpha < 0.999:
            a = sp_pr.makeelement(qn("a:alpha"), {"val": str(int(max(0.0, alpha) * 100000))})
            clr.append(a)
        return clr

    def color_el(parent_tag: str, color: str | None, alpha: float) -> object:
        wrap = sp_pr.makeelement(qn(parent_tag), {})
        wrap.append(srgb_el(color, alpha))
        return wrap

    if grad.get("kind") == "pattern":
        ang = grad.get("angle", 45) % 180
        prst = min(((0, "horz"), (45, "dnDiag"), (90, "vert"), (135, "upDiag"), (180, "horz")),
                   key=lambda kv: abs(kv[0] - ang))[1]
        patt = sp_pr.makeelement(qn("a:pattFill"), {"prst": prst})
        patt.append(color_el("a:fgClr", grad.get("fg"), 1.0))
        patt.append(color_el("a:bgClr", grad.get("bg") or grad.get("fg"), 1.0))
        sp_pr.append(patt)
        return

    stops = sorted(grad.get("stops") or [], key=lambda s: s.get("f", 0))
    if len(stops) < 2:
        return
    # Transparent stops have no color of their own — borrow the nearest
    # colored neighbor so the fade doesn't pass through gray.
    colored = [s for s in stops if s.get("color")]
    fallback = colored[0]["color"] if colored else "FFFFFF"
    grad_fill = sp_pr.makeelement(qn("a:gradFill"), {"rotWithShape": "1"})
    gs_lst = sp_pr.makeelement(qn("a:gsLst"), {})
    for s in stops:
        pos = int(min(1.0, max(0.0, s.get("f", 0))) * 100000)
        gs = sp_pr.makeelement(qn("a:gs"), {"pos": str(pos)})
        gs.append(srgb_el(s.get("color") or fallback, s.get("alpha", 1.0)))
        gs_lst.append(gs)
    grad_fill.append(gs_lst)
    if grad.get("kind") == "radial":
        at = grad.get("at") or {"x": 0.5, "y": 0.5}
        path = sp_pr.makeelement(qn("a:path"), {"path": "circle"})
        fr = sp_pr.makeelement(qn("a:fillToRect"), {
            "l": str(int(at["x"] * 100000)), "t": str(int(at["y"] * 100000)),
            "r": str(int((1 - at["x"]) * 100000)), "b": str(int((1 - at["y"]) * 100000)),
        })
        path.append(fr)
        grad_fill.append(path)
    else:
        ang = int(((grad.get("angle", 180) - 90) % 360) * 60000)
        lin = sp_pr.makeelement(qn("a:lin"), {"ang": str(ang), "scaled": "1"})
        grad_fill.append(lin)
    sp_pr.append(grad_fill)


def strip_theme_style(shp) -> None:
    """Remove the <p:style> theme reference python-pptx puts on autoshapes.

    Its effectRef idx="2" points at the default Office theme's effect
    style, which carries an outer shadow — so every rect renders with a
    phantom drop shadow even though the shape itself declares none. We
    set explicit fills/lines anyway, so the whole reference is noise.
    """
    el = shp._element
    style = el.find(qn("p:style"))
    if style is not None:
        el.remove(style)


def emit_pptx(slides: list, out_path: Path, slide_w: int, slide_h: int) -> None:
    prs = Presentation()
    prs.slide_width  = Emu(px_to_emu(slide_w))
    prs.slide_height = Emu(px_to_emu(slide_h))
    blank = prs.slide_layouts[6]

    for s in slides:
        slide = prs.slides.add_slide(blank)
        # Background rect spanning the whole slide.
        if s.get("bg"):
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                prs.slide_width, prs.slide_height)
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = hex_to_rgb(s["bg"])
            bg_shape.line.fill.background()
            strip_theme_style(bg_shape)

        for p in s.get("primitives", []):
            x = px_to_emu(p["x"])
            y = px_to_emu(p["y"])
            w = max(px_to_emu(p["w"]), 1)
            h = max(px_to_emu(p["h"]), 1)
            kind = p["kind"]

            if kind == "rect":
                radius = p.get("radius", 0)
                corner = p.get("corner")
                if corner in ("top", "bottom", "left", "right") and radius >= 3:
                    # Rect with one rounded corner-pair (header cells,
                    # accent strips, first/last bar segments). The
                    # shape rounds its TOP pair; rotation points that
                    # pair at the right edge. left/right need the
                    # geometry swapped so the rotated bounding box
                    # matches the target rect.
                    shape_type = MSO_SHAPE.ROUND_2_SAME_RECTANGLE
                elif radius >= 4:
                    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                else:
                    shape_type = MSO_SHAPE.RECTANGLE
                if shape_type == MSO_SHAPE.ROUND_2_SAME_RECTANGLE and corner in ("left", "right"):
                    cx, cy = x + w // 2, y + h // 2
                    w, h = h, w
                    x, y = cx - w // 2, cy - h // 2
                shp = slide.shapes.add_shape(shape_type, x, y, w, h)
                if shape_type == MSO_SHAPE.ROUND_2_SAME_RECTANGLE:
                    short = min(p["w"], p["h"])
                    if short > 0 and shp.adjustments:
                        shp.adjustments[0] = min(0.5, radius / short)
                        if len(shp.adjustments) > 1:
                            shp.adjustments[1] = 0
                    shp.rotation = {"top": 0, "right": 90, "bottom": 180, "left": 270}[corner]
                elif shape_type == MSO_SHAPE.ROUNDED_RECTANGLE and shp.adjustments:
                    # Adjustment 0 is 0..0.5 of short side; convert px → fraction.
                    short = min(p["w"], p["h"])
                    if short > 0:
                        shp.adjustments[0] = min(0.5, radius / short)
                if p.get("grad"):
                    apply_gradient_fill(shp, p["grad"])
                elif p.get("fill"):
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = hex_to_rgb(p["fill"])
                else:
                    shp.fill.background()
                if p.get("line") and p.get("line_w", 0) > 0.5:
                    shp.line.color.rgb = hex_to_rgb(p["line"])
                    shp.line.width = Pt(max(0.5, p["line_w"] * 0.75))
                else:
                    shp.line.fill.background()
                strip_theme_style(shp)

            elif kind == "text":
                tb = slide.shapes.add_textbox(x, y, w, h)
                if p.get("rot"):
                    tb.rotation = float(p["rot"])
                tf = tb.text_frame
                tf.word_wrap = False
                tf.margin_left = tf.margin_right = Emu(0)
                tf.margin_top = tf.margin_bottom = Emu(0)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE if p.get("anchor") == "middle" else MSO_ANCHOR.TOP
                if p.get("vert"):
                    # writing-mode: vertical-rl → native rotated text flow.
                    tf._txBody.bodyPr.set("vert", p["vert"])
                para = tf.paragraphs[0]
                para.space_before = Pt(0)
                para.space_after  = Pt(0)
                para.line_spacing = 1.0
                align = p.get("align") or "left"
                para.alignment = {
                    "left":   PP_ALIGN.LEFT,
                    "right":  PP_ALIGN.RIGHT,
                    "center": PP_ALIGN.CENTER,
                    "justify": PP_ALIGN.JUSTIFY,
                }.get(align, PP_ALIGN.LEFT)
                # Each primitive carries N styled runs (one per
                # style transition within the line). Bold/italic/
                # colored inline spans land here as separate runs
                # sharing the same paragraph + position.
                runs = p.get("runs") or [p]   # back-compat with flat shape
                for r_data in runs:
                    run = para.add_run()
                    run.text = r_data.get("text", "")
                    run.font.name = r_data.get("font_family", "Helvetica")
                    pt_size = max(6.0, r_data["font_size_px"] * 0.75)
                    run.font.size = Pt(pt_size)
                    if r_data.get("font_weight", 400) >= 600:
                        run.font.bold = True
                    if r_data.get("font_style") == "italic":
                        run.font.italic = True
                    if r_data.get("underline"):
                        run.font.underline = True
                    run.font.color.rgb = hex_to_rgb(r_data.get("color", "000000"))
                    # Kern above 12pt like browsers do — without this,
                    # large display numerals render with uneven glyph
                    # gaps (PowerPoint/Keynote default kerning off).
                    run.font._rPr.set("kern", "1200")
                    ls_px = r_data.get("letter_spacing_px", 0) or 0
                    if abs(ls_px) > 0.01:
                        run.font._rPr.set("spc", str(int(round(ls_px * 75))))
                    if r_data.get("strikethrough"):
                        run.font._rPr.set("strike", "sngStrike")

            elif kind == "pic":
                blob = fetch_image(p["src"])
                if not blob:
                    continue
                try:
                    pic = slide.shapes.add_picture(
                        io.BytesIO(blob), x, y, width=w, height=h)
                except Exception as e:
                    print(f"  warn: pic insert failed: {e}", file=sys.stderr)

    prs.save(str(out_path))


def main() -> int:
    ap = argparse.ArgumentParser(prog="pptx", description="DOM-walking HTML → PPTX exporter")
    ap.add_argument("--url", required=True, help="URL of the deck.html (must be reachable from headless chromium)")
    ap.add_argument("--out", required=True, help="output .pptx path")
    ap.add_argument("--width",  type=int, default=1920)
    ap.add_argument("--height", type=int, default=1080)
    args = ap.parse_args()

    out_path = Path(args.out).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # Import playwright here so the script can `--help` without it.
    from playwright.sync_api import sync_playwright

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": args.width, "height": args.height})
        page.goto(args.url, wait_until="networkidle", timeout=30000)
        page.wait_for_function("customElements.get('deck-stage') !== undefined", timeout=15000)
        # Settle a tick to let the shell mount its DOM.
        page.wait_for_timeout(400)
        slides = page.evaluate(DOM_WALK_JS)
        browser.close()

    if not slides:
        print("no slides found — is <deck-stage> present?", file=sys.stderr)
        return 1
    emit_pptx(slides, out_path, args.width, args.height)
    print(f"wrote {out_path}  ({len(slides)} slide(s))")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
